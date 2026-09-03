"""Slide PDF ingestion.

Text is kept as one string with `[Slide N]` markers rather than a row per page,
so the summariser can cite a slide number and nothing about the database schema
has to change. Scanned, image-only decks yield no text — pypdf does not do OCR,
and that is out of scope, so those slides are marked rather than silently empty.
"""

from __future__ import annotations

import re
from pathlib import Path

# PDF text extraction produces ragged whitespace: trailing spaces from column
# layouts, and runs of blank lines between text boxes.
_TRAILING_SPACE = re.compile(r"[ \t]+$", re.MULTILINE)
_BLANK_RUN = re.compile(r"\n{3,}")

NO_TEXT_MARKER = "(no extractable text — this slide is probably an image)"


def _tidy(text: str) -> str:
    return _BLANK_RUN.sub("\n\n", _TRAILING_SPACE.sub("", text)).strip()


def extract_pdf_text(pdf_path: str | Path) -> tuple[str, int]:
    """Return (text with [Slide N] markers, page count)."""
    from pypdf import PdfReader

    reader = PdfReader(str(pdf_path))

    if reader.is_encrypted:
        # Many university decks are "protected" with an empty owner password,
        # which opens fine; a real password is the user's to remove.
        try:
            opened = reader.decrypt("")
        except Exception:
            opened = 0
        if not opened:
            raise ValueError("This PDF is password-protected — remove the password and re-upload")

    pages: list[str] = []
    for number, page in enumerate(reader.pages, start=1):
        try:
            body = _tidy(page.extract_text() or "")
        except Exception:
            body = ""  # one malformed page must not lose the rest of the deck
        pages.append(f"[Slide {number}]\n{body or NO_TEXT_MARKER}")

    return "\n\n".join(pages), len(reader.pages)


def page_count(extracted_text: str) -> int:
    """Slides in an already-extracted deck, without reopening the file."""
    return len(re.findall(r"^\[Slide \d+\]$", extracted_text, re.MULTILINE))


def has_text(extracted_text: str) -> bool:
    """False when every page came back empty — i.e. the deck needs OCR to be
    useful, so the UI can say so instead of implying the slides were read."""
    stripped = re.sub(r"^\[Slide \d+\]$", "", extracted_text, flags=re.MULTILINE)
    return bool(stripped.replace(NO_TEXT_MARKER, "").strip())


def _shape_text(shape) -> list[str]:
    """Text from one shape, recursing into groups and tables.

    Lecturers put content in all three, and a deck that loses its tables loses
    exactly the dense reference material worth studying from.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    parts: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            parts.extend(_shape_text(child))
        return parts

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
        return parts

    if getattr(shape, "has_text_frame", False):
        body = shape.text_frame.text.strip()
        if body:
            parts.append(body)
    return parts


def extract_pptx_text(pptx_path: str | Path) -> tuple[str, int]:
    """Return (text with [Slide N] markers, slide count) read straight from the
    .pptx.

    Deliberately not done by converting to PDF first: a PDF export drops the
    speaker notes, which are often the most useful part of a lecturer's deck —
    the things they meant to say rather than what fit on the slide.
    """
    from pptx import Presentation

    presentation = Presentation(str(pptx_path))

    slides: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_text(shape))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker notes: {notes}")

        body = _tidy("\n".join(parts))
        slides.append(f"[Slide {number}]\n{body or NO_TEXT_MARKER}")

    return "\n\n".join(slides), len(presentation.slides)


# Suffixes we can pull text out of directly, and how.
_EXTRACTORS = {".pdf": extract_pdf_text, ".pptx": extract_pptx_text}

# Legacy binary PowerPoint. python-pptx cannot read it, so it has to be
# converted to PDF first and is only accepted when a converter exists.
NEEDS_CONVERSION = {".ppt"}

SUPPORTED_SUFFIXES = set(_EXTRACTORS) | NEEDS_CONVERSION


def extract_text(path: str | Path) -> tuple[str, int]:
    """Dispatch on file type. Raises ValueError for anything unsupported."""
    suffix = Path(path).suffix.lower()
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise ValueError(f"Cannot read text from a '{suffix}' file")
    return extractor(path)
